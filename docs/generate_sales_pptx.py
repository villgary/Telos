#!/usr/bin/env python3
"""
从 A2-sales-deck-outline.md 生成 Telos 销售 PPTX
改进版：slide编号、进度条、表格自适应、摘要高亮框、emoji处理
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re, os, sys

# ── 颜色定义 ──────────────────────────────────────────────
DARK_BLUE    = RGBColor(0x0D, 0x2B, 0x5A)
MED_BLUE     = RGBColor(0x1E, 0x5A, 0x9E)
ACCENT_BLUE  = RGBColor(0x21, 0x8A, 0xD4)
LIGHT_BLUE   = RGBColor(0xE8, 0xF4, 0xFD)
RED_ACCENT   = RGBColor(0xE5, 0x3E, 0x3E)
ORANGE       = RGBColor(0xF5, 0xA6, 0x23)
GREEN        = RGBColor(0x52, 0xC4, 0x1A)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT    = RGBColor(0x1F, 0x1F, 0x2E)
LIGHT_GRAY   = RGBColor(0xF5, 0xF7, 0xFA)
MID_GRAY     = RGBColor(0x89, 0x97, 0xA8)
HIGHLIGHT_BG = RGBColor(0xFF, 0xF3, 0xE0)   # 淡橙色高亮背景
SECTION_BG   = RGBColor(0x0A, 0x1F, 0x4A)   # 深色章节分隔页背景

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── 工具函数 ──────────────────────────────────────────────

def clean(text):
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)
    return text.strip()


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=14, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_footer(slide, page_num, total):
    """页面底部：品牌 + 页码"""
    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.4), LIGHT_GRAY)
    add_text(slide, "Telos  |  身份威胁检测与响应平台",
             Inches(0.4), Inches(7.15), Inches(6), Inches(0.3),
             size=9, color=MID_GRAY)
    add_text(slide, f"{page_num} / {total}",
             Inches(12.0), Inches(7.15), Inches(1.0), Inches(0.3),
             size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def calc_col_widths(headers, rows, total_width_inches, min_col=1.5, max_col=6.0):
    """根据内容长度自适应计算列宽"""
    n = len(headers)
    # 收集每列内容长度（字符数）
    col_lens = [len(h) for h in headers]
    for row in rows:
        for i, cell in enumerate(row):
            if i < n:
                col_lens[i] = max(col_lens[i], len(clean(str(cell))))

    total_len = sum(max(cl, 1) for cl in col_lens)
    # 分配比例权重
    raw_widths = [total_width_inches * cl / total_len for cl in col_lens]
    # 裁剪到 [min_col, max_col] 范围
    clamped = [max(min_col, min(max_col, rw)) for rw in raw_widths]
    # 补偿溢出（如果总和超过 total_width）
    excess = sum(clamped) - total_width_inches
    if excess > 0:
        scale = total_width_inches / sum(clamped)
        clamped = [max(min_col, w * scale) for w in clamped]
    return clamped


def add_table(slide, headers, rows, left, top, width, height, col_widths=None):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths is None:
        col_widths = calc_col_widths(headers, rows, float(width / 914400 * 96))

    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw)

    for i, h in enumerate(headers):
        c = tbl.cell(0, i)
        c.text = clean(h)
        c.fill.solid()
        c.fill.fore_color.rgb = DARK_BLUE
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0] if p.runs else p.add_run()
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE

    for ri, row in enumerate(rows):
        bg = LIGHT_GRAY if ri % 2 else WHITE
        for ci, ct in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            c.text = clean(str(ct))
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]
            r = p.runs[0] if p.runs else p.add_run()
            r.font.size = Pt(10)
            r.font.color.rgb = DARK_TEXT


# ── 幻灯片生成 ─────────────────────────────────────────────

def make_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    # 顶部装饰线
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), ACCENT_BLUE)
    add_rect(slide, 0, Inches(0.12), SLIDE_W, Inches(0.05), ORANGE)

    # 主标题
    add_text(slide, "TELOS",
             Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2),
             size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 英文 slogan
    add_text(slide, "Know every account. Detect every threat. Respond before damage.",
             Inches(0.8), Inches(3.0), Inches(11.7), Inches(0.8),
             size=22, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # 中文副标题
    for i, line in enumerate([
            "身份威胁检测与响应平台",
            "Identity Threat Detection & Response"
    ]):
        color = RGBColor(0xCC, 0xD6, 0xE8) if i == 0 else ACCENT_BLUE
        size = 18 if i == 0 else 15
        add_text(slide, line,
                 Inches(0.8), Inches(4.0) + Inches(i * 0.58),
                 Inches(11.7), Inches(0.55),
                 size=size, color=color, align=PP_ALIGN.CENTER, italic=(i == 1))

    # 底部
    add_rect(slide, 0, Inches(6.7), SLIDE_W, Inches(0.8), MED_BLUE)
    add_text(slide, "www.telos.com  |  contact@telos.com  |  400-XXX-XXXX",
             Inches(0), Inches(6.85), SLIDE_W, Inches(0.5),
             size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)


def make_toc(prs, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    # 左侧色条
    add_rect(slide, 0, 0, Inches(0.4), SLIDE_H, DARK_BLUE)
    add_rect(slide, Inches(0.4), 0, Inches(0.08), SLIDE_H, ACCENT_BLUE)

    add_text(slide, "目录", Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             size=36, bold=True, color=DARK_BLUE)
    add_text(slide, "TABLE OF CONTENTS",
             Inches(0.8), Inches(1.05), Inches(5), Inches(0.35),
             size=12, color=MID_GRAY, italic=True)

    n_cols, per_col = 2, 7
    for idx, (num, cn, en) in enumerate(items):
        col, row = idx // per_col, idx % per_col
        lx = Inches(0.8 + col * 6.2)
        ty = Inches(1.6 + row * 0.68)
        add_rect(slide, lx, ty + Inches(0.04), Inches(0.28), Inches(0.28), MED_BLUE)
        add_text(slide, str(num), lx, ty + Inches(0.02), Inches(0.28), Inches(0.32),
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, cn, lx + Inches(0.38), ty, Inches(5.5), Inches(0.35),
                 size=13, color=DARK_TEXT)
        add_text(slide, en, lx + Inches(0.38), ty + Inches(0.28), Inches(5.5), Inches(0.3),
                 size=9, color=MID_GRAY, italic=True)


def make_content(prs, num, title_cn, title_en, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.08), DARK_BLUE)
    add_rect(slide, 0, Inches(1.08), SLIDE_W, Inches(0.06), ORANGE)
    # 页码标记
    add_rect(slide, Inches(12.7), Inches(0.28), Inches(0.48), Inches(0.48), ORANGE)
    add_text(slide, str(num), Inches(12.7), Inches(0.28), Inches(0.48), Inches(0.48),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title_cn, Inches(0.5), Inches(0.18), Inches(12.0), Inches(0.52),
             size=24, bold=True, color=WHITE)
    add_text(slide, title_en, Inches(0.5), Inches(0.65), Inches(12.0), Inches(0.35),
             size=13, color=ACCENT_BLUE, italic=True)
    add_footer(slide, page_num, total)
    return slide


def make_section_divider(prs, section_num, section_cn, section_en, page_num, total):
    """章节分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SECTION_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), ACCENT_BLUE)
    add_rect(slide, 0, Inches(0.12), SLIDE_W, Inches(0.05), ORANGE)
    add_text(slide, f"0{section_num}",
             Inches(0.8), Inches(1.5), Inches(2.0), Inches(1.5),
             size=96, bold=True, color=MED_BLUE, align=PP_ALIGN.LEFT)
    add_text(slide, section_cn,
             Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.0),
             size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, section_en,
             Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.6),
             size=20, color=ACCENT_BLUE, align=PP_ALIGN.LEFT, italic=True)
    add_footer(slide, page_num, total)


def add_bullets(slide, items, top=Inches(1.35), left=Inches(0.4), w=Inches(12.5),
                font_size=12, line_h=Inches(0.50)):
    """渲染要点列表，支持 emoji 前缀"""
    ty = top
    for item in items:
        text = clean(item)
        icon = ""

        # emoji 检测（更宽松： Unicode > 0x2600 或已知的符号前缀）
        if text and ord(text[0]) >= 0x2600:
            icon = text[0]
            text = text[1:].lstrip()
        elif text.startswith("✅") or text.startswith("❌") or \
             text.startswith("⚠️") or text.startswith("🔴") or \
             text.startswith("🟠") or text.startswith("🟡") or \
             text.startswith("🟢") or text.startswith("📊") or \
             text.startswith("🔗") or text.startswith("🤖") or \
             text.startswith("📤") or text.startswith("⚡"):
            icon = text[0]
            text = text[1:].lstrip()
        elif len(text) >= 2 and text[:2] in ["V ", "X ", "! "]:
            icon = text[:2]
            text = text[2:].lstrip()
        elif text.startswith("+") or text.startswith("-"):
            icon = text[0]
            text = text[1:].lstrip()

        display = (icon + "  " + text) if icon else text
        add_text(slide, display,
                 left, ty, w, line_h,
                 size=font_size, color=DARK_TEXT)
        ty += line_h


def add_highlight_box(slide, text, icon="💡", top=Inches(5.8)):
    """高亮摘要框（橙色底）"""
    add_rect(slide, Inches(0.4), top, Inches(12.5), Inches(0.6),
             HIGHLIGHT_BG, line_color=ORANGE, line_width=Pt(1.5))
    add_rect(slide, Inches(0.4), top, Inches(0.06), Inches(0.6), ORANGE)
    add_text(slide, f"{icon}  {clean(text)}",
             Inches(0.6), top + Inches(0.08),
             Inches(12.0), Inches(0.48),
             size=11, color=DARK_TEXT, italic=True)


def add_quote(slide, text, top=Inches(6.3)):
    add_rect(slide, Inches(0.4), top, Inches(12.5), Inches(0.55), LIGHT_BLUE,
             line_color=ACCENT_BLUE, line_width=Pt(1))
    add_text(slide, "\u201c" + clean(text) + "\u201d",
             Inches(0.6), top + Inches(0.08),
             Inches(12.0), Inches(0.42),
             size=11, color=DARK_BLUE, italic=True)


def add_screenshot_placeholder(slide, label_cn, label_en, top=Inches(1.3),
                               height=Inches(3.5)):
    """产品截图占位符"""
    add_rect(slide, Inches(0.4), top, Inches(7.8), height,
             LIGHT_GRAY, line_color=MID_GRAY, line_width=Pt(0.75))
    add_text(slide, f"[ {label_cn} ]",
             Inches(0.4), top + height / 2 - Inches(0.3),
             Inches(7.8), Inches(0.3),
             size=14, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, f"[ {label_en} ]",
             Inches(0.4), top + height / 2,
             Inches(7.8), Inches(0.3),
             size=11, color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)


# ── Markdown 解析 ───────────────────────────────────────────

SECTION_MAP = {
    3: (1, "市场背景", "Market Context"),
    5: (2, "核心痛点", "Core Pain Points"),
    6: (3, "Telos 解决方案", "Telos Solution"),
    7: (4, "核心能力", "Core Capabilities"),
    14: (5, "产品界面", "Product UI"),
    17: (6, "客户价值", "Customer Value"),
    19: (7, "竞争优势", "Competitive Edge"),
    21: (8, "客户案例", "Customer Stories"),
    23: (9, "部署与生态", "Deployment & Ecosystem"),
}


def parse_md(path):
    with open(path, encoding="utf-8") as f:
        text = f.read()
    slides = []
    text = re.sub(r'^---\n.*?\n---\n', '', text, count=1, flags=re.DOTALL)
    for block in re.split(r'\n---\n', text):
        block = block.strip()
        if not block:
            continue
        m = re.match(r'##\s+Slide\s+(\d+)\s*[-—]\s*(.+)', block)
        if not m:
            continue
        num = int(m.group(1))
        full_title = m.group(2).strip()
        parts = full_title.split("/")
        cn = parts[0].strip()
        en = parts[1].strip() if len(parts) > 1 else ""

        data = {
            "num": num,
            "cn": cn,
            "en": en,
            "block": block,
        }

        # 提取引言
        quote_m = re.search(
            r'\*\*(?:导言|Quote|Key Message|核心信息)[：:]*\*\*\s*>?\s*[""]?(.+)',
            block, re.DOTALL)
        if quote_m:
            data["quote"] = clean(quote_m.group(1))

        # 提取引言块（多行）
        quote_block_m = re.search(r'^>\s*(.+)$', block, re.MULTILINE)
        if quote_block_m and "quote" not in data:
            data["quote"] = clean(quote_block_m.group(1))

        # 提取高亮摘要（**亮点 / Highlight:**）
        hl_m = re.search(
            r'\*\*(?:亮点|Highlight|核心亮点)[：:]*\*\*\s*(.+)',
            block, re.DOTALL)
        if hl_m:
            data["highlight"] = clean(hl_m.group(1))

        # 提取架构文本
        arch_m = re.search(r'```[^\n]*\n(.+?)```', block, re.DOTALL)
        if arch_m:
            data["arch"] = arch_m.group(1).strip()

        # 提取截图占位符标记 <!-- screenshot: xxx / xxx -->
        ss_m = re.search(r'<!--\s*screenshot:\s*(.+?)\s*/\s*(.+?)\s*-->', block)
        if ss_m:
            data["screenshot_cn"] = ss_m.group(1).strip()
            data["screenshot_en"] = ss_m.group(2).strip()

        # 提取表格
        tlines = [l for l in block.split("\n")
                  if "|" in l and "---" not in l and l.strip().startswith("|")]
        if len(tlines) >= 2:
            headers = [h.strip() for h in tlines[0].split("|") if h.strip()]
            rows = []
            for rl in tlines[2:]:
                cells = [c.strip() for c in rl.split("|")[1:-1]]
                if cells:
                    rows.append(cells)
            if rows and len(headers) >= 2:
                data["headers"] = headers
                data["rows"] = rows

        # 提取要点
        bullets = []
        for line in block.split("\n"):
            line = line.strip()
            if re.match(r'^[-*+]\s', line):
                bullets.append(clean(line[2:]))
            elif re.match(r'^\d+[.、\)]\s', line):
                bullets.append(clean(re.sub(r'^\d+[.、\)]\s*', '', line)))
        data["bullets"] = bullets

        # 判断 slide 类型
        if "screenshot_cn" in data:
            slides.append(("screenshot", data))
        elif "headers" in data:
            slides.append(("table", data))
        elif "arch" in data:
            slides.append(("arch", data))
        else:
            slides.append(("content", data))

    return slides


# ── 主函数 ────────────────────────────────────────────────

def main():
    docs_dir = os.path.dirname(os.path.abspath(__file__))
    md_path = os.path.join(docs_dir, "A2-sales-deck-outline.md")
    out_path = os.path.join(docs_dir, "A2-sales-deck-Telos.pptx")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 删除默认空白页
    if len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    slide_list = parse_md(md_path)

    # ── 计算总页数（封面 + 目录 + 内容页 + 章节分隔 + 结束页）
    total_pages = 2  # 封面 + 目录
    shown_nums = set()
    for stype, sdata in slide_list:
        num = sdata["num"]
        if num in (1, 2, 25):
            continue
        shown_nums.add(num)
        # 章节分隔页
        if num in SECTION_MAP:
            total_pages += 1
        total_pages += 1
    total_pages += 1  # 结束页

    # 进度打印
    def progress(msg):
        print(msg, flush=True)

    progress(f"[1/{total_pages}] 生成封面...")
    make_cover(prs)

    # 解析目录
    toc_items = []
    for stype, sdata in slide_list:
        if sdata["num"] == 2:
            block = sdata.get("block", "")
            for line in block.split("\n"):
                m = re.match(r'^\d+[.、]\s*(.+)', line.strip())
                if m:
                    title = m.group(1).strip()
                    parts = title.split("/")
                    cn = parts[0].strip() if parts else title
                    en = parts[1].strip() if len(parts) > 1 else ""
                    toc_items.append((len(toc_items) + 1, cn, en))
    progress(f"[2/{total_pages}] 生成目录...")
    make_toc(prs, toc_items or [(i+1, f"内容 {i+1}", "") for i in range(11)])

    # 内容页
    page_counter = 2
    shown_nums_sorted = sorted(shown_nums)
    num_to_page = {n: i + 3 for i, n in enumerate(shown_nums_sorted)}
    # 重新计算含章节分隔页的页码
    page_counter = 2
    for num in shown_nums_sorted:
        if num in SECTION_MAP:
            page_counter += 1
        page_counter += 1
    real_total = page_counter + 1  # +1 结束页

    page_counter = 2
    for num in shown_nums_sorted:
        section_info = SECTION_MAP.get(num)
        if section_info:
            page_counter += 1
            progress(f"[{page_counter}/{real_total}] 章节分隔：{section_info[1]}...")
            make_section_divider(prs, section_info[0], section_info[1],
                                 section_info[2], page_counter, real_total)

        page_counter += 1

        # 找到对应的 slide 数据
        sdata = next(
            (sd for stype_, sd in slide_list if sd["num"] == num), None)
        if not sdata:
            continue

        stype = next(
            (st for st, sd in slide_list if sd["num"] == num), "")
        cn, en = sdata["cn"], sdata["en"]
        progress(f"[{page_counter}/{real_total}] Slide {num}: {cn}")

        slide = make_content(prs, num, cn, en, page_counter, real_total)

        if stype == "screenshot":
            # 左侧：截图占位，右侧：要点
            sc_cn = sdata.get("screenshot_cn", "产品截图")
            sc_en = sdata.get("screenshot_en", "Product Screenshot")
            add_screenshot_placeholder(slide, sc_cn, sc_en,
                                       top=Inches(1.3), height=Inches(3.2))
            bullets = sdata.get("bullets", [])
            if bullets:
                # 右侧要点区
                add_bullets(slide, bullets,
                            top=Inches(1.3), left=Inches(8.4),
                            w=Inches(4.5), font_size=11, line_h=Inches(0.48))
            hl = sdata.get("highlight", "")
            if hl:
                add_highlight_box(slide, hl, icon="💡", top=Inches(5.0))

        elif stype == "table":
            headers = sdata["headers"]
            rows = sdata["rows"]
            n = len(headers)
            tbl_top = Inches(1.3)
            tbl_h = Inches(min(len(rows) + 1, 7) * 0.48)
            col_widths = calc_col_widths(headers, rows, 12.5)
            add_table(slide, headers, rows,
                      Inches(0.4), tbl_top, Inches(12.5), tbl_h,
                      col_widths=col_widths)
            bullets = sdata.get("bullets", [])
            if bullets:
                bullet_top = tbl_top + tbl_h + Inches(0.1)
                add_bullets(slide, bullets, top=bullet_top,
                             font_size=11, line_h=Inches(0.46))
            hl = sdata.get("highlight", "")
            if hl:
                add_highlight_box(slide, hl, icon="💡",
                                  top=tbl_top + tbl_h + Inches(0.55))

        elif stype == "arch":
            arch = sdata.get("arch", "")
            if arch:
                add_rect(slide, Inches(0.4), Inches(1.3),
                         Inches(12.5), Inches(3.8), LIGHT_GRAY,
                         line_color=MID_GRAY, line_width=Pt(0.5))
                add_text(slide, arch,
                         Inches(0.6), Inches(1.4),
                         Inches(12.1), Inches(3.6),
                         size=10, color=DARK_BLUE)
            bullets = sdata.get("bullets", [])
            if bullets:
                bullet_top = Inches(1.3) + (Inches(3.8) if arch else Inches(0))
                add_bullets(slide, bullets, top=bullet_top,
                             font_size=11, line_h=Inches(0.46))
            hl = sdata.get("highlight", "")
            if hl:
                add_highlight_box(slide, hl, icon="💡", top=Inches(5.0))
            quote = sdata.get("quote", "")
            if quote:
                add_quote(slide, quote, top=Inches(6.1))

        else:
            bullets = sdata.get("bullets", [])
            if bullets:
                add_bullets(slide, bullets, font_size=12, line_h=Inches(0.50))
            hl = sdata.get("highlight", "")
            if hl:
                add_highlight_box(slide, hl, icon="💡", top=Inches(5.8))
            quote = sdata.get("quote", "")
            if quote:
                add_quote(slide, quote, top=Inches(6.3))

    progress(f"[{real_total}/{real_total}] 生成结束页...")
    # 结束页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.1), ACCENT_BLUE)
    add_rect(slide, 0, Inches(0.1), SLIDE_W, Inches(0.05), ORANGE)
    add_text(slide, "立即开启身份安全新篇章",
             0, Inches(2.4), SLIDE_W, Inches(1.0),
             size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Start Your Identity Security Journey Today",
             0, Inches(3.4), SLIDE_W, Inches(0.6),
             size=20, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, italic=True)
    for i, line in enumerate([
            "www.telos.com  |  contact@telos.com  |  400-XXX-XXXX",
            "",
            "预约产品演示  |  申请 POC 环境  |  安全评估咨询"
    ]):
        add_text(slide, line, 0, Inches(4.3) + Inches(i * 0.5),
                 SLIDE_W, Inches(0.5),
                 size=14, color=RGBColor(0xCC, 0xD6, 0xE8),
                 align=PP_ALIGN.CENTER)
    add_rect(slide, 0, Inches(6.8), SLIDE_W, Inches(0.7), MED_BLUE)
    add_text(slide, "Telos v2.0  |  © 2026 Telos",
             0, Inches(6.9), SLIDE_W, Inches(0.5),
             size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

    prs.save(out_path)
    print(f"\n✅ 完成：{out_path} （共 {len(prs.slides)} 页）")


if __name__ == "__main__":
    main()
